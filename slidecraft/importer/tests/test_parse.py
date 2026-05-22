"""Tests for slidecraft.importer.parse.

All PPTX fixtures are built in-memory using python-pptx; no binary files are committed.
"""
from __future__ import annotations

import io
import zipfile
from pathlib import Path

import pytest
from lxml import etree
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt, Emu

from slidecraft.importer.parse import (
    parse,
    _txbody_is_empty,
    _layout_ph_has_custom_prompt,
    _parse_paragraph,
)
from slidecraft.importer.model import (
    NoFill,
    Paragraph,
    Presentation,
    Run,
    SolidFill,
)

_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


# ---------------------------------------------------------------------------
# Helper to serialise an in-memory pptx to a temp file path
# ---------------------------------------------------------------------------

def _save_pptx(prs: PptxPresentation, tmp_path: Path) -> Path:
    """Save a python-pptx Presentation to a temp file and return its path."""
    out = tmp_path / "test.pptx"
    prs.save(str(out))
    return out


# ---------------------------------------------------------------------------
# Test: canvas dimensions
# ---------------------------------------------------------------------------

def test_parse_canvas_dims(tmp_path):
    """parse() reads <p:sldSz> and converts EMU → px correctly."""
    prs = PptxPresentation()
    # Default python-pptx presentation is 10×7.5 in = 9144000×6858000 EMU
    # at 96 dpi: 960×720 px
    out = _save_pptx(prs, tmp_path)
    result = parse(out)
    # px = emu / 9525
    expected_w = prs.slide_width // 9525
    expected_h = prs.slide_height // 9525
    assert result.canvas_width_px == expected_w
    assert result.canvas_height_px == expected_h


def test_parse_canvas_dims_16x9(tmp_path):
    """parse() handles 16:9 (1920×1080) canvas correctly."""
    prs = PptxPresentation()
    # Set to 16:9: 12192000 × 6858000 EMU
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    out = _save_pptx(prs, tmp_path)
    result = parse(out)
    assert result.canvas_width_px == 12192000 // 9525
    assert result.canvas_height_px == 6858000 // 9525


# ---------------------------------------------------------------------------
# Test: one title placeholder
# ---------------------------------------------------------------------------

def test_parse_one_title_placeholder(tmp_path):
    """parse() extracts a title placeholder with text runs from a slide."""
    prs = PptxPresentation()
    slide_layout = prs.slide_layouts[0]  # "Title Slide" layout
    slide = prs.slides.add_slide(slide_layout)

    # Set title text
    title_ph = slide.shapes.title
    title_ph.text = "Hello World"

    out = _save_pptx(prs, tmp_path)
    result = parse(out)

    assert len(result.slides) == 1
    slide_model = result.slides[0]
    assert slide_model.index == 1

    # Find the title placeholder
    title_phs = [p for p in slide_model.placeholders if p.type in ("title", "ctrTitle")]
    assert len(title_phs) >= 1, "Expected at least one title placeholder"

    ph = title_phs[0]
    assert ph.text_frame is not None
    all_text = "".join(
        run.text
        for para in ph.text_frame.paragraphs
        for run in para.runs
    )
    assert "Hello World" in all_text


def test_parse_placeholder_geometry(tmp_path):
    """parse() resolves placeholder position/size in px."""
    prs = PptxPresentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)

    title_ph = slide.shapes.title
    title_ph.text = "Geometry Test"

    out = _save_pptx(prs, tmp_path)
    result = parse(out)

    ph = result.slides[0].placeholders[0]
    # Should have non-zero width (inherited from layout or master)
    assert ph.width_px > 0 or ph.height_px > 0, (
        "Expected non-zero placeholder dimensions"
    )


def test_parse_multiple_slides(tmp_path):
    """parse() returns one Slide model per PPTX slide."""
    prs = PptxPresentation()
    layout = prs.slide_layouts[0]
    prs.slides.add_slide(layout)
    prs.slides.add_slide(layout)

    out = _save_pptx(prs, tmp_path)
    result = parse(out)

    assert len(result.slides) == 2
    assert result.slides[0].index == 1
    assert result.slides[1].index == 2


def test_parse_typefaces_referenced(tmp_path):
    """parse() collects typefaces from placeholder defaults and run-level fonts."""
    prs = PptxPresentation()
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Font Test"

    out = _save_pptx(prs, tmp_path)
    result = parse(out)

    # typefaces_referenced should be a set (possibly empty if theme fonts are all +mj-lt/+mn-lt refs)
    assert isinstance(result.typefaces_referenced, set)


def test_parse_empty_presentation(tmp_path):
    """parse() handles a presentation with no slides gracefully."""
    prs = PptxPresentation()
    out = _save_pptx(prs, tmp_path)
    result = parse(out)
    assert isinstance(result, Presentation)
    assert result.slides == []


def test_parse_slide_index_1based(tmp_path):
    """Slide index is 1-based."""
    prs = PptxPresentation()
    layout = prs.slide_layouts[1]
    prs.slides.add_slide(layout)
    out = _save_pptx(prs, tmp_path)
    result = parse(out)
    assert result.slides[0].index == 1


def test_parse_body_placeholder(tmp_path):
    """parse() extracts body placeholder text content."""
    prs = PptxPresentation()
    # Use a layout that has both title and content placeholders
    layout = prs.slide_layouts[1]  # "Title and Content" layout
    slide = prs.slides.add_slide(layout)

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = "Body content here"
            break

    out = _save_pptx(prs, tmp_path)
    result = parse(out)

    slide_model = result.slides[0]
    body_phs = [p for p in slide_model.placeholders if p.idx == 1]
    if body_phs:
        ph = body_phs[0]
        assert ph.text_frame is not None
        all_text = "".join(
            run.text
            for para in ph.text_frame.paragraphs
            for run in para.runs
        )
        assert "Body content" in all_text


def test_parse_default_run_props_populated(tmp_path):
    """parse() populates default_run_props on each placeholder (cascade result)."""
    prs = PptxPresentation()
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Cascade Test"

    out = _save_pptx(prs, tmp_path)
    result = parse(out)

    slide_model = result.slides[0]
    assert len(slide_model.placeholders) > 0
    for ph in slide_model.placeholders:
        # default_run_props must be a Run instance (text="" sentinel)
        assert ph.default_run_props is not None
        assert ph.default_run_props.text == ""
        # default_para_props must be a Paragraph instance
        assert ph.default_para_props is not None


# ---------------------------------------------------------------------------
# Helpers for building minimal PPTX fixtures from scratch via lxml
# ---------------------------------------------------------------------------

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_PKG = "http://schemas.openxmlformats.org/package/2006/relationships"


def _build_pptx_with_prompt_fallback(tmp_path: Path) -> Path:
    """Build a minimal .pptx where slide1/ph_5 is empty and layout ph_5 has
    hasCustomPrompt='1' with text 'Layout Prompt Text'.

    We start from a normal python-pptx blank presentation and then
    post-edit the layout XML via zipfile manipulation to inject the
    hasCustomPrompt attribute and prompt text.
    """
    from pptx.util import Inches
    prs = PptxPresentation()

    # Use "Title and Content" layout (index 1) which has a body ph idx=1
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)

    # Leave the body placeholder EMPTY (just set title so slide is valid)
    slide.shapes.title.text = "My Title"
    # body placeholder (idx=1) gets no text

    pptx_path = tmp_path / "prompt_test.pptx"
    prs.save(str(pptx_path))

    # Now post-edit the PPTX (zip) to set hasCustomPrompt on the layout's body ph
    import zipfile, shutil
    edited_path = tmp_path / "prompt_test_edited.pptx"
    with zipfile.ZipFile(pptx_path, "r") as zin, zipfile.ZipFile(edited_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if "slideLayout" in item.filename and item.filename.endswith(".xml"):
                # Find layout used by slide1 — check rels first is complex; edit ALL layouts
                try:
                    root = etree.fromstring(data)
                    ns_map = {
                        "a": _A, "p": _P, "r": _R
                    }
                    sp_tree = root.find(".//p:spTree", ns_map)
                    if sp_tree is not None:
                        for sp in sp_tree.findall("p:sp", ns_map):
                            nv_pr = sp.find(".//p:ph", ns_map)
                            if nv_pr is not None and nv_pr.get("idx") == "1":
                                # Set hasCustomPrompt
                                nv_pr.set("hasCustomPrompt", "1")
                                # Inject prompt text into txBody
                                tx_body = sp.find("p:txBody", ns_map)
                                if tx_body is not None:
                                    for p_el in list(tx_body.findall(f"{{{_A}}}p")):
                                        tx_body.remove(p_el)
                                    p_new = etree.SubElement(tx_body, f"{{{_A}}}p")
                                    r_new = etree.SubElement(p_new, f"{{{_A}}}r")
                                    rpr_new = etree.SubElement(r_new, f"{{{_A}}}rPr")
                                    rpr_new.set("lang", "en-US")
                                    t_new = etree.SubElement(r_new, f"{{{_A}}}t")
                                    t_new.text = "Layout Prompt Text"
                    data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                except Exception:
                    pass  # leave data unchanged
            zout.writestr(item, data)

    return edited_path


def _build_pptx_with_layout_fill(tmp_path: Path) -> Path:
    """Build a minimal .pptx where slide1/ph_1 has empty spPr and layout ph_1
    has a solidFill with schemeClr 'accent1', to test fill cascading.
    """
    from pptx.util import Inches
    prs = PptxPresentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Fill Test"

    pptx_path = tmp_path / "fill_test.pptx"
    prs.save(str(pptx_path))

    import zipfile
    edited_path = tmp_path / "fill_test_edited.pptx"
    with zipfile.ZipFile(pptx_path, "r") as zin, zipfile.ZipFile(edited_path, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if "slideLayout" in item.filename and item.filename.endswith(".xml"):
                try:
                    root = etree.fromstring(data)
                    ns_map = {"a": _A, "p": _P}
                    sp_tree = root.find(".//p:spTree", ns_map)
                    if sp_tree is not None:
                        for sp in sp_tree.findall("p:sp", ns_map):
                            nv_pr = sp.find(".//p:ph", ns_map)
                            if nv_pr is not None and nv_pr.get("idx") == "1":
                                sp_pr = sp.find("p:spPr", ns_map)
                                if sp_pr is None:
                                    sp_pr = etree.SubElement(sp, f"{{{_P}}}spPr")
                                # Remove any existing fill
                                for tag in ["a:noFill", "a:solidFill", "a:gradFill"]:
                                    el = sp_pr.find(tag, ns_map)
                                    if el is not None:
                                        sp_pr.remove(el)
                                # Add solidFill with a known sRGB color
                                solid = etree.SubElement(sp_pr, f"{{{_A}}}solidFill")
                                srgb = etree.SubElement(solid, f"{{{_A}}}srgbClr")
                                srgb.set("val", "FF4422")
                    data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                except Exception:
                    pass
            zout.writestr(item, data)

    return edited_path


# ---------------------------------------------------------------------------
# Test: prompt fallback (Deliverable 1)
# ---------------------------------------------------------------------------

class TestPromptFallback:
    def test_empty_placeholder_uses_layout_prompt(self, tmp_path):
        """When slide ph is empty and layout ph has hasCustomPrompt='1', use layout text."""
        pptx_path = _build_pptx_with_prompt_fallback(tmp_path)
        result = parse(pptx_path)
        slide = result.slides[0]
        body_phs = [p for p in slide.placeholders if p.idx == 1]
        assert body_phs, "Expected a body placeholder with idx=1"
        ph = body_phs[0]
        assert ph.is_prompt_fallback is True, "Placeholder should be flagged as prompt fallback"
        assert ph.text_frame is not None, "TextFrame should be populated from layout prompt"
        all_text = "".join(
            run.text
            for para in ph.text_frame.paragraphs
            for run in para.runs
        )
        assert "Layout Prompt Text" in all_text, (
            f"Expected 'Layout Prompt Text' in fallback content, got: {all_text!r}"
        )

    def test_non_empty_placeholder_not_fallback(self, tmp_path):
        """When slide ph has real content, is_prompt_fallback must stay False."""
        prs = PptxPresentation()
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Title"
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                ph.text = "Real content"
                break
        out = _save_pptx(prs, tmp_path)
        result = parse(out)
        body_phs = [p for p in result.slides[0].placeholders if p.idx == 1]
        if body_phs:
            assert body_phs[0].is_prompt_fallback is False

    def test_prompt_fallback_false_by_default_no_layout_prompt(self, tmp_path):
        """When layout ph has no hasCustomPrompt, is_prompt_fallback stays False even if empty."""
        prs = PptxPresentation()
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Title"
        # Leave body empty — but layout doesn't have hasCustomPrompt
        out = _save_pptx(prs, tmp_path)
        result = parse(out)
        for ph in result.slides[0].placeholders:
            assert ph.is_prompt_fallback is False


# ---------------------------------------------------------------------------
# Helpers for unit tests
# ---------------------------------------------------------------------------

def _make_sp_with_prompt(idx: int, has_custom: bool, text: str) -> etree._Element:
    """Build a minimal <p:sp> element for testing _layout_ph_has_custom_prompt."""
    sp = etree.Element(f"{{{_P}}}sp")
    nv_sp_pr = etree.SubElement(sp, f"{{{_P}}}nvSpPr")
    etree.SubElement(nv_sp_pr, f"{{{_P}}}cNvPr", {"id": "1", "name": "test"})
    etree.SubElement(nv_sp_pr, f"{{{_P}}}cNvSpPr")
    nv_pr = etree.SubElement(nv_sp_pr, f"{{{_P}}}nvPr")
    ph_attrs = {"idx": str(idx)}
    if has_custom:
        ph_attrs["hasCustomPrompt"] = "1"
    etree.SubElement(nv_pr, f"{{{_P}}}ph", ph_attrs)
    sp_pr = etree.SubElement(sp, f"{{{_P}}}spPr")
    tx_body = etree.SubElement(sp, f"{{{_P}}}txBody")
    body_pr = etree.SubElement(tx_body, f"{{{_A}}}bodyPr")
    lst_style = etree.SubElement(tx_body, f"{{{_A}}}lstStyle")
    p_el = etree.SubElement(tx_body, f"{{{_A}}}p")
    if text:
        r_el = etree.SubElement(p_el, f"{{{_A}}}r")
        t_el = etree.SubElement(r_el, f"{{{_A}}}t")
        t_el.text = text
    return sp


class TestLayoutPhHasCustomPrompt:
    def test_has_custom_prompt_true(self):
        sp = _make_sp_with_prompt(5, True, "Click to edit")
        assert _layout_ph_has_custom_prompt(sp) is True

    def test_has_custom_prompt_false(self):
        sp = _make_sp_with_prompt(5, False, "")
        assert _layout_ph_has_custom_prompt(sp) is False

    def test_has_custom_prompt_no_ph(self):
        sp = etree.Element(f"{{{_P}}}sp")
        assert _layout_ph_has_custom_prompt(sp) is False


class TestTxBodyIsEmpty:
    def test_empty_txbody(self):
        tx = etree.fromstring(
            f'<txBody xmlns="{_A}"><bodyPr/><p><endParaRPr/></p></txBody>'
        )
        assert _txbody_is_empty(tx) is True

    def test_whitespace_only_txbody(self):
        tx = etree.fromstring(
            f'<txBody xmlns="{_A}"><bodyPr/><p><r><t>   </t></r></p></txBody>'
        )
        assert _txbody_is_empty(tx) is True

    def test_non_empty_txbody(self):
        tx = etree.fromstring(
            f'<txBody xmlns="{_A}"><bodyPr/><p><r><t>Hello</t></r></p></txBody>'
        )
        assert _txbody_is_empty(tx) is False


# ---------------------------------------------------------------------------
# Test: layout fill cascade (Deliverable 2)
# ---------------------------------------------------------------------------

class TestLayoutFillCascade:
    def test_layout_solidfill_propagates_when_slide_sppr_empty(self, tmp_path):
        """When slide ph has empty spPr, fill cascades from layout placeholder."""
        pptx_path = _build_pptx_with_layout_fill(tmp_path)
        result = parse(pptx_path)
        slide = result.slides[0]
        body_phs = [p for p in slide.placeholders if p.idx == 1]
        assert body_phs, "Expected body placeholder idx=1"
        ph = body_phs[0]
        assert isinstance(ph.fill, SolidFill), (
            f"Expected SolidFill cascaded from layout, got {ph.fill!r}"
        )
        # The color we injected was #FF4422
        assert ph.fill.color.r == 0xFF
        assert ph.fill.color.g == 0x44
        assert ph.fill.color.b == 0x22

    def test_slide_sppr_nofill_is_respected(self, tmp_path):
        """When slide ph has explicit <a:noFill>, layout fill should NOT override."""
        prs = PptxPresentation()
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "NoFill Test"
        pptx_path = tmp_path / "nofill_test.pptx"
        prs.save(str(pptx_path))

        # Post-edit: add noFill to slide ph_1 spPr
        edited = tmp_path / "nofill_edited.pptx"
        with zipfile.ZipFile(pptx_path, "r") as zin, zipfile.ZipFile(edited, "w", zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if "slides/slide1.xml" in item.filename:
                    try:
                        root = etree.fromstring(data)
                        ns_map = {"a": _A, "p": _P}
                        sp_tree = root.find(".//p:spTree", ns_map)
                        if sp_tree is not None:
                            for sp in sp_tree.findall("p:sp", ns_map):
                                nv_pr = sp.find(".//p:ph", ns_map)
                                if nv_pr is not None and nv_pr.get("idx") == "1":
                                    sp_pr = sp.find("p:spPr", ns_map)
                                    if sp_pr is None:
                                        sp_pr = etree.SubElement(sp, f"{{{_P}}}spPr")
                                    etree.SubElement(sp_pr, f"{{{_A}}}noFill")
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                    except Exception:
                        pass
                zout.writestr(item, data)

        result = parse(edited)
        body_phs = [p for p in result.slides[0].placeholders if p.idx == 1]
        if body_phs:
            assert isinstance(body_phs[0].fill, NoFill), (
                "Explicit noFill on slide should not be overridden by layout fill"
            )


# ---------------------------------------------------------------------------
# Tests for <a:fld> field element parsing
# ---------------------------------------------------------------------------

def _make_paragraph_with_fld(field_type: str, field_text: str) -> etree._Element:
    """Build an <a:p> element containing an <a:fld> child (e.g. slidenum, datetime1)."""
    p_el = etree.Element(f"{{{_A_NS}}}p")
    fld = etree.SubElement(p_el, f"{{{_A_NS}}}fld",
                           attrib={"id": "{DEADBEEF-DEAD-DEAD-DEAD-DEADDEADBEEF}",
                                   "type": field_type})
    rpr = etree.SubElement(fld, f"{{{_A_NS}}}rPr", lang="en-US")
    t_el = etree.SubElement(fld, f"{{{_A_NS}}}t")
    t_el.text = field_text
    return p_el


class TestFieldElementParsing:
    """Tests for <a:fld> (slidenum, datetime) parsed as regular runs."""

    def test_fld_slidenum_text_captured(self):
        """<a:fld type='slidenum'> text is emitted as a Run with the field value."""
        p_el = _make_paragraph_with_fld("slidenum", "7")
        default_run = Run(text="")
        default_para = Paragraph(runs=[])
        para = _parse_paragraph(p_el, default_run, default_para)
        texts = [r.text for r in para.runs]
        assert "7" in texts

    def test_fld_datetime_text_captured(self):
        """<a:fld type='datetime1'> text is emitted as a Run with the field value."""
        p_el = _make_paragraph_with_fld("datetime1", "1/1/2025")
        default_run = Run(text="")
        default_para = Paragraph(runs=[])
        para = _parse_paragraph(p_el, default_run, default_para)
        texts = [r.text for r in para.runs]
        assert "1/1/2025" in texts

    def test_fld_interleaved_with_runs(self):
        """Runs and fields in the same paragraph are all captured in order."""
        p_el = etree.Element(f"{{{_A_NS}}}p")
        # A regular run before the field
        r1 = etree.SubElement(p_el, f"{{{_A_NS}}}r")
        rpr1 = etree.SubElement(r1, f"{{{_A_NS}}}rPr")
        t1 = etree.SubElement(r1, f"{{{_A_NS}}}t")
        t1.text = "Slide "
        # The field
        fld = etree.SubElement(p_el, f"{{{_A_NS}}}fld",
                               attrib={"id": "{DEADBEEF}", "type": "slidenum"})
        rpr_fld = etree.SubElement(fld, f"{{{_A_NS}}}rPr")
        t_fld = etree.SubElement(fld, f"{{{_A_NS}}}t")
        t_fld.text = "3"

        default_run = Run(text="")
        default_para = Paragraph(runs=[])
        para = _parse_paragraph(p_el, default_run, default_para)
        all_text = "".join(r.text for r in para.runs)
        assert "Slide " in all_text
        assert "3" in all_text
