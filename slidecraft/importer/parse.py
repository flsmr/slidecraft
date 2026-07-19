"""PPTX parser — Stage 0 of the PPTX → Slidev pipeline.

parse(pptx_path) → Presentation

Parses text-bearing placeholders only (v1 scope).  Drops to lxml for properties
that python-pptx doesn't surface (bodyPr insets, normAutofit fontScale, line/para
spacing, bullet chars, etc.).
"""
from __future__ import annotations

import re
from pathlib import Path
from typing import Optional

from lxml import etree
import pptx
from pptx import Presentation as PptxPresentation
from pptx.util import Emu

from slidecraft.importer.model import (
    Fill,
    GradientStop,
    LinearGradientFill,
    NoFill,
    Paragraph,
    Picture,
    Placeholder,
    Presentation,
    RGB,
    RadialGradientFill,
    Run,
    Slide,
    SolidFill,
    TextFrame,
)
from slidecraft.importer.inheritance import (
    _extract_rpr,
    _extract_ppr,
    _find,
    _get,
    _parse_color,
    _x,
    diff_run,
    diff_para,
    get_effective_clr_map,
    resolve_placeholder,
)
from slidecraft.importer.pictures.effects import parse_effects
from slidecraft.importer.pictures.geometry import cust_geom_to_clip_path, preset_to_css
# walk_text_shapes is imported lazily inside parse() — shapes.parse imports
# private helpers from THIS module, so a top-level import would deadlock.

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
_EMU_PER_PX = 9525          # 914400 EMU/inch ÷ 96 px/inch
_EMU_PER_PT = 12700         # 914400 EMU/inch ÷ 72 pt/inch

# Placeholder types considered text-bearing
_TEXT_PH_TYPES = {
    "title", "ctrTitle", "body", "subTitle", "dt", "ftr", "sldNum", None
}

# ---------------------------------------------------------------------------
# EMU converters
# ---------------------------------------------------------------------------

def _emu_to_px(emu: int) -> float:
    return emu / _EMU_PER_PX


def _emu_to_pt(emu: int) -> float:
    return emu / _EMU_PER_PT


# ---------------------------------------------------------------------------
# Fill resolver
# ---------------------------------------------------------------------------

def _resolve_fill(
    sp_el: etree._Element,
    theme_el: Optional[etree._Element],
    clr_map: Optional[dict[str, str]] = None,
    layout_sp: Optional[etree._Element] = None,
) -> Fill:
    """Resolve placeholder fill using a slide → layout cascade.

    For placeholder shapes the OOXML spec says that if the slide-level <p:spPr>
    contains no fill directive (no <a:solidFill>, <a:gradFill>, or <a:noFill>),
    the fill should be inherited from the corresponding layout placeholder.  We
    implement that two-level cascade here so that layout-defined solid fills
    (like the ``bg2`` band boxes on IU slide 1) appear correctly.
    """
    def _read_fill_from_sp_pr(sp_pr: Optional[etree._Element]) -> Optional[Fill]:
        """Return a Fill or None if sp_pr carries no fill directive."""
        if sp_pr is None:
            return None
        if sp_pr.find(_x("a:noFill")) is not None:
            return NoFill()
        solid = sp_pr.find(_x("a:solidFill"))
        if solid is not None:
            color = _parse_color(solid, theme_el, clr_map)
            if color:
                return SolidFill(color)
        grad = sp_pr.find(_x("a:gradFill"))
        if grad is not None:
            return _parse_grad_fill(grad, theme_el, clr_map)
        return None  # no directive found

    # Level 1: slide-level shape
    slide_sp_pr = sp_el.find(_x("p:spPr"))
    fill = _read_fill_from_sp_pr(slide_sp_pr)
    if fill is not None:
        return fill

    # Level 2: layout placeholder (cascade fallback)
    if layout_sp is not None:
        layout_sp_pr = layout_sp.find(_x("p:spPr"))
        fill = _read_fill_from_sp_pr(layout_sp_pr)
        if fill is not None:
            return fill

    return NoFill()


def _parse_grad_fill(
    grad_el: etree._Element,
    theme_el: Optional[etree._Element] = None,
    clr_map: Optional[dict[str, str]] = None,
) -> Fill:
    """Parse <a:gradFill> into LinearGradientFill or RadialGradientFill."""
    stops: list[GradientStop] = []
    gs_lst = grad_el.find(_x("a:gsLst"))
    if gs_lst is not None:
        for gs in gs_lst.findall(_x("a:gs")):
            pos_str = gs.get("pos", "0")
            pos = int(pos_str) / 100000.0
            solid = gs.find(_x("a:solidFill"))
            color = _parse_color(solid, theme_el, clr_map) if solid is not None else RGB(0, 0, 0)
            if color:
                stops.append(GradientStop(position=pos, color=color))

    # Linear vs radial
    lin = grad_el.find(_x("a:lin"))
    path_el = grad_el.find(_x("a:path"))
    if lin is not None:
        ang_str = lin.get("ang", "0")
        # PPT angle is in 60000ths of a degree, measured clockwise from the top
        angle_deg = int(ang_str) / 60000.0
        return LinearGradientFill(angle_deg=angle_deg, stops=stops)
    elif path_el is not None and path_el.get("path") == "circle":
        return RadialGradientFill(stops=stops)
    else:
        return LinearGradientFill(angle_deg=0.0, stops=stops)


# ---------------------------------------------------------------------------
# Slide background resolver
# ---------------------------------------------------------------------------

def _resolve_background(
    slide_part,
    pptx_prs: PptxPresentation,
    theme_el: Optional[etree._Element],
    clr_map: Optional[dict[str, str]] = None,
) -> Fill:
    """Walk slide → layout → master for <p:bg> / <p:bgRef> to find the background fill."""
    for part in [slide_part, slide_part.slide_layout, slide_part.slide_layout.slide_master]:
        sp_tree = part._element.find(_x("p:cSld"))
        if sp_tree is None:
            continue
        bg = sp_tree.find(_x("p:bg"))
        if bg is None:
            continue
        bg_pr = bg.find(_x("p:bgPr"))
        if bg_pr is not None:
            no_fill = bg_pr.find(_x("a:noFill"))
            if no_fill is not None:
                return NoFill()
            solid = bg_pr.find(_x("a:solidFill"))
            if solid is not None:
                color = _parse_color(solid, theme_el, clr_map)
                if color:
                    return SolidFill(color)
            grad = bg_pr.find(_x("a:gradFill"))
            if grad is not None:
                return _parse_grad_fill(grad, theme_el, clr_map)
        # bgRef references a fill style in the theme — not resolved in v1
        bg_ref = bg.find(_x("p:bgRef"))
        if bg_ref is not None:
            solid = bg_ref.find(_x("a:solidFill"))
            if solid is None:
                # srgbClr directly under bgRef (unusual but valid)
                pass
            else:
                color = _parse_color(solid, theme_el, clr_map)
                if color:
                    return SolidFill(color)
    return NoFill()


# ---------------------------------------------------------------------------
# TextFrame parser
# ---------------------------------------------------------------------------

def _parse_text_frame(
    tx_body: etree._Element,
    default_run: Run,
    default_para: Paragraph,
    theme_el: Optional[etree._Element] = None,
    clr_map: Optional[dict[str, str]] = None,
    layout_tx_body: Optional[etree._Element] = None,
    master_tx_body: Optional[etree._Element] = None,
    per_level_defaults: Optional[dict[int, tuple[Run, Paragraph]]] = None,
) -> TextFrame:
    """Parse <p:txBody> into a TextFrame, diffing each run/paragraph against the defaults.

    bodyPr attributes (anchor, insets, rotation, autofit) cascade slide → layout
    → master per OOXML semantics: a missing or unset attribute on a lower-priority
    bodyPr inherits from the next level up. The IU template's date/footer
    placeholders rely on the master cascade (slide has empty <a:bodyPr/>, layout
    doesn't redeclare the placeholder, master sets lIns=tIns=rIns=bIns=0).
    Hardcoded OOXML defaults apply only when none of slide/layout/master set
    the attribute.
    """
    # Collect bodyPr from each cascade level (any may be None).
    slide_body_pr = tx_body.find(_x("a:bodyPr"))
    layout_body_pr = (
        layout_tx_body.find(_x("a:bodyPr")) if layout_tx_body is not None else None
    )
    master_body_pr = (
        master_tx_body.find(_x("a:bodyPr")) if master_tx_body is not None else None
    )

    def _attr_cascade(name: str) -> Optional[str]:
        """Return the first non-None occurrence of bodyPr@name across slide → layout → master."""
        for src in (slide_body_pr, layout_body_pr, master_body_pr):
            if src is not None:
                v = src.get(name)
                if v is not None:
                    return v
        return None

    def _child_cascade(local_name: str) -> Optional[etree._Element]:
        """Return the first child element with the given local name across slide → layout → master."""
        for src in (slide_body_pr, layout_body_pr, master_body_pr):
            if src is not None:
                el = src.find(_x(f"a:{local_name}"))
                if el is not None:
                    return el
        return None

    anchor_val = _attr_cascade("anchor") or "t"
    anchor = anchor_val  # "t", "ctr", "b"

    l_ins = _attr_cascade("lIns")
    t_ins = _attr_cascade("tIns")
    r_ins = _attr_cascade("rIns")
    b_ins = _attr_cascade("bIns")
    # Hardcoded OOXML defaults: lIns=91440, tIns=45720, rIns=91440, bIns=45720
    l_pt = _emu_to_pt(int(l_ins)) if l_ins is not None else 7.2
    t_pt = _emu_to_pt(int(t_ins)) if t_ins is not None else 3.6
    r_pt = _emu_to_pt(int(r_ins)) if r_ins is not None else 7.2
    b_pt = _emu_to_pt(int(b_ins)) if b_ins is not None else 3.6
    insets_pt = (l_pt, t_pt, r_pt, b_pt)

    rot = _attr_cascade("rot")
    rotation_deg = int(rot) / 60000.0 if rot is not None else 0.0

    norm_auto = _child_cascade("normAutofit")
    autofit_font_scale = 1.0
    if norm_auto is not None:
        fs_str = norm_auto.get("fontScale", "100000")
        autofit_font_scale = int(fs_str) / 100000.0

    # Parse paragraphs.  Each <a:p> may expand into multiple Paragraphs if it
    # contains a hard line break (\r\n or \n) embedded inside an <a:t> — PPT
    # treats those as paragraph breaks (new bullet), not soft line breaks.
    paragraphs: list[Paragraph] = []
    for p_el in tx_body.findall(_x("a:p")):
        paragraphs.extend(
            _parse_paragraph(
                p_el, default_run, default_para, theme_el, clr_map,
                per_level_defaults=per_level_defaults,
            )
        )

    return TextFrame(
        paragraphs=paragraphs,
        anchor=anchor,  # type: ignore[arg-type]
        insets_pt=insets_pt,
        rotation_deg=rotation_deg,
        autofit_font_scale=autofit_font_scale,
    )


_RUN_BACKFILL_FIELDS = (
    "font_size_pt", "bold", "italic", "underline", "strike",
    "color", "font_family", "cap",
)


class _ParsedParagraphs(list):
    """Result of :func:`_parse_paragraph` — a ``list[Paragraph]`` with one
    entry per hard-break segment, plus a ``runs`` view spanning them all.

    ``runs`` returns every Run parsed from the original ``<a:p>`` in
    document order, so callers can treat the result as "the paragraph's
    content" without caring whether the hard-break split produced one
    segment or several.
    """

    @property
    def runs(self) -> list[Run]:
        return [r for p in self for r in p.runs]


def _parse_paragraph(
    p_el: etree._Element,
    default_run: Run,
    default_para: Paragraph,
    theme_el: Optional[etree._Element] = None,
    clr_map: Optional[dict[str, str]] = None,
    per_level_defaults: Optional[dict[int, tuple[Run, Paragraph]]] = None,
) -> "_ParsedParagraphs":
    """Parse one ``<a:p>`` into one OR MORE :class:`Paragraph` objects.

    A single PPT paragraph may expand into multiple model Paragraphs when
    a hard line break (``\\r\\n`` / ``\\n`` / ``\\r``) appears INSIDE one
    of its ``<a:t>`` runs. PowerPoint treats such embedded breaks as
    paragraph separators (each segment becomes a NEW bullet when the
    paragraph is bulleted) — see the IU template's slide 5 body, where
    a single ``<a:p>`` ends one run with ``"...17 pt.\\r\\n"`` followed
    by the next run starting "Um …" and PowerPoint renders that as two
    bullets. This is distinct from the soft-break ``<a:br/>`` element,
    which stays within a paragraph.

    All produced Paragraphs share the same pPr / level / diff defaults
    (cascaded from the input ``<a:p>``).
    """
    ppr_el = p_el.find(_x("a:pPr"))
    para_props = _extract_ppr(ppr_el, theme_el, clr_map)

    # Level (indent level, 0-based)
    level = 0
    if ppr_el is not None:
        lvl_str = ppr_el.get("lvl", "0")
        level = int(lvl_str)

    # Build the diff'd paragraph TEMPLATE — we deep-clone this per split.
    diffed_template = diff_para(para_props, default_para)
    diffed_template.level = level

    # Walk the <a:p>'s direct children in document order, accumulating runs
    # into the CURRENT paragraph segment. A ``\\r\\n`` inside an <a:t>
    # closes the current segment and opens a new one (paragraph break);
    # ``<a:br/>`` stays within the current segment as a Run(text="\\n")
    # marker which the emit layer renders as <br/>.
    segments: list[list[Run]] = [[]]   # at least one segment

    # Per-level cascade backfill. When the paragraph's level > 0 AND we
    # have per-level resolved defaults, any RUN field the slide didn't set
    # explicitly is backfilled with the level-N default value BEFORE the
    # diff against the placeholder's level-0 default. This makes diff_run
    # surface deviations like `font_size_pt=20` for lvl=3 runs (master's
    # lvl4pPr.defRPr.sz), even when the slide-level <a:rPr> omits sz.
    #
    # Without this, lvl=N paragraphs inherited the placeholder wrapper's
    # level-0 font-size, producing the visible "sub-level bullets render
    # at the SAME large size as level-0" symptom on tmp2 slide 10.
    level_run: Optional[Run] = None
    if level > 0 and per_level_defaults is not None:
        entry = per_level_defaults.get(level)
        if entry is not None:
            level_run = entry[0]

    def _make_diffed_run(text: str, rpr_el: Optional[etree._Element]) -> Run:
        run_props = _extract_rpr(rpr_el, theme_el, clr_map)
        if level_run is not None:
            for f in _RUN_BACKFILL_FIELDS:
                if getattr(run_props, f) is None:
                    setattr(run_props, f, getattr(level_run, f))
        run_props.text = text
        diffed_run = diff_run(run_props, default_run)
        diffed_run.text = text
        return diffed_run

    def _add_text_with_paragraph_breaks(
        text: str,
        rpr_el: Optional[etree._Element],
    ) -> None:
        """Split *text* on universal newlines; each break closes the current
        paragraph segment and opens a new one carrying the SAME pPr."""
        # Normalise CRLF / CR → LF so we can split on a single delimiter.
        normalised = text.replace("\r\n", "\n").replace("\r", "\n")
        parts = normalised.split("\n")
        for i, part in enumerate(parts):
            if i > 0:
                # Close current segment, open a new one.
                segments.append([])
            # Always emit a run (even empty) so explicit rPr stays attached.
            segments[-1].append(_make_diffed_run(part, rpr_el))

    for child in p_el:
        tag = etree.QName(child.tag).localname
        if tag == "r":
            t_el = child.find(_x("a:t"))
            text = t_el.text or "" if t_el is not None else ""
            rpr_el = child.find(_x("a:rPr"))
            _add_text_with_paragraph_breaks(text, rpr_el)
        elif tag == "fld":
            # Field element: <a:fld type="slidenum"|"datetime1"|…>
            t_el = child.find(_x("a:t"))
            text = t_el.text or "" if t_el is not None else ""
            rpr_el = child.find(_x("a:rPr"))
            _add_text_with_paragraph_breaks(text, rpr_el)
        elif tag == "br":
            # Soft line break — stays inside the current paragraph segment.
            segments[-1].append(Run(text="\n"))

    # Materialise one Paragraph per segment, all sharing the template's pPr.
    out = _ParsedParagraphs()
    for seg in segments:
        p = Paragraph(
            runs=seg,
            align=diffed_template.align,
            line_spacing_pct=diffed_template.line_spacing_pct,
            space_before_pt=diffed_template.space_before_pt,
            space_after_pt=diffed_template.space_after_pt,
            indent_pt=diffed_template.indent_pt,
            margin_left_pt=diffed_template.margin_left_pt,
            bullet=diffed_template.bullet,
            bullet_char=diffed_template.bullet_char,
            bullet_color=diffed_template.bullet_color,
            bullet_font=diffed_template.bullet_font,
            bullet_size_pct=diffed_template.bullet_size_pct,
            bullet_autonum_type=diffed_template.bullet_autonum_type,
            level=diffed_template.level,
        )
        out.append(p)
    return out


# ---------------------------------------------------------------------------
# Placeholder locator helpers
# ---------------------------------------------------------------------------

def _ph_idx(sp_el: etree._Element) -> Optional[int]:
    """Return the idx of a <p:sp>'s placeholder, or None if not a placeholder."""
    nv_sp_pr = sp_el.find(_x("p:nvSpPr"))
    if nv_sp_pr is None:
        return None
    nv_pr = nv_sp_pr.find(_x("p:nvPr"))
    if nv_pr is None:
        return None
    ph = nv_pr.find(_x("p:ph"))
    if ph is None:
        return None
    return int(ph.get("idx", "0"))


def _ph_type(sp_el: etree._Element) -> Optional[str]:
    """Return the type attribute of a placeholder (may be None for body/generic)."""
    nv_sp_pr = sp_el.find(_x("p:nvSpPr"))
    if nv_sp_pr is None:
        return None
    nv_pr = nv_sp_pr.find(_x("p:nvPr"))
    if nv_pr is None:
        return None
    ph = nv_pr.find(_x("p:ph"))
    if ph is None:
        return None
    return ph.get("type")  # e.g. "title", "body", None


def _parse_body_pr_autofit(sp_el: etree._Element) -> tuple[Optional[bool], Optional[bool]]:
    """Read autofit + wrap from a shape's ``<a:bodyPr>``.

    Returns ``(shape_autofit, wrap_text)`` where each may be ``None`` if the
    setting is absent at this level — the caller is responsible for
    cascading slide → layout → master via ``or``-like resolution.

    OOXML semantics:
      - ``<a:spAutoFit/>`` child   → ``shape_autofit=True``  (box grows to fit text)
      - ``<a:normAutofit/>`` child → ``shape_autofit=False`` (text shrinks instead)
      - ``<a:noAutofit/>`` child   → ``shape_autofit=False`` (default behaviour)
      - bodyPr present but no autofit child → ``shape_autofit=None`` (inherit)
      - ``bodyPr.@wrap="none"``    → ``wrap_text=False`` (extend horizontally,
                                      never line-wrap)
      - ``bodyPr.@wrap="square"`` or absent → ``wrap_text=True`` (default)

    The combination ``wrap_text=False`` + ``shape_autofit=True`` is what
    PowerPoint generates for "title-style" placeholders that the designer
    wants to size to their content (e.g. the IU template's slide-24
    centre title at width 540px doesn't wrap "Designs with tables" — the
    box is expected to extend horizontally to fit the actual content).
    """
    tx_body = sp_el.find(_x("p:txBody"))
    if tx_body is None:
        return None, None
    body_pr = tx_body.find(_x("a:bodyPr"))
    if body_pr is None:
        return None, None

    # Autofit: examine the bodyPr's children for the three known markers.
    shape_autofit: Optional[bool] = None
    if body_pr.find(_x("a:spAutoFit")) is not None:
        shape_autofit = True
    elif body_pr.find(_x("a:noAutofit")) is not None:
        shape_autofit = False
    elif body_pr.find(_x("a:normAutofit")) is not None:
        # normAutofit means "scale text down to fit"; we don't currently
        # implement font-scaling on the emit side, so semantically the
        # SHAPE is fixed — shape_autofit=False expresses that.
        shape_autofit = False

    # Wrap: bodyPr.@wrap defaults to "square" (wrap-enabled). We only flip
    # to False when it's explicitly "none".
    wrap_attr = body_pr.get("wrap")
    if wrap_attr == "none":
        wrap_text: Optional[bool] = False
    elif wrap_attr is not None:
        wrap_text = True
    else:
        wrap_text = None  # inherit

    return shape_autofit, wrap_text


def _resolve_body_pr_cascade(
    slide_sp: Optional[etree._Element],
    layout_sp: Optional[etree._Element],
    master_sp: Optional[etree._Element],
) -> tuple[bool, bool]:
    """Cascade slide → layout → master to resolve bodyPr autofit/wrap.

    Returns ``(shape_autofit, wrap_text)`` with concrete bools (no None).
    Cascade order matches OOXML: the slide's own setting wins if present,
    else inherit the layout's, else the master's, else PPT defaults
    (``shape_autofit=False, wrap_text=True``).

    Each property cascades INDEPENDENTLY — a slide might set wrap=none
    explicitly while leaving autofit to inherit from the layout.
    """
    shape_autofit: Optional[bool] = None
    wrap_text: Optional[bool] = None
    for src in (slide_sp, layout_sp, master_sp):
        if src is None:
            continue
        sa, wt = _parse_body_pr_autofit(src)
        if shape_autofit is None and sa is not None:
            shape_autofit = sa
        if wrap_text is None and wt is not None:
            wrap_text = wt
        if shape_autofit is not None and wrap_text is not None:
            break  # both resolved, no need to descend further
    return (
        shape_autofit if shape_autofit is not None else False,
        wrap_text if wrap_text is not None else True,
    )


def _ph_has_text(sp_el: etree._Element) -> bool:
    """Return True if the <p:sp> has a <p:txBody> with at least one <a:t> with text."""
    tx_body = sp_el.find(_x("p:txBody"))
    if tx_body is None:
        return False
    for t_el in tx_body.iter(_x("a:t")):
        if t_el.text:
            return True
    return False


def _txbody_is_empty(tx_body: etree._Element) -> bool:
    """Return True if all <a:t> runs in a txBody are missing or contain only whitespace."""
    for t_el in tx_body.iter(_x("a:t")):
        if t_el.text and t_el.text.strip():
            return False
    return True


def _layout_ph_has_custom_prompt(layout_sp: etree._Element) -> bool:
    """Return True if the layout <p:sp> has hasCustomPrompt='1' on its <p:ph>."""
    nv_sp_pr = layout_sp.find(_x("p:nvSpPr"))
    if nv_sp_pr is None:
        return False
    nv_pr = nv_sp_pr.find(_x("p:nvPr"))
    if nv_pr is None:
        return False
    ph = nv_pr.find(_x("p:ph"))
    if ph is None:
        return False
    return ph.get("hasCustomPrompt") == "1"


def _get_sp_position(sp_el: etree._Element) -> tuple[float, float, float, float, float]:
    """Return (x_px, y_px, width_px, height_px, rotation_deg) from <p:spPr><a:xfrm>."""
    sp_pr = sp_el.find(_x("p:spPr"))
    if sp_pr is None:
        return 0.0, 0.0, 0.0, 0.0, 0.0
    xfrm = sp_pr.find(_x("a:xfrm"))
    if xfrm is None:
        return 0.0, 0.0, 0.0, 0.0, 0.0
    off = xfrm.find(_x("a:off"))
    ext = xfrm.find(_x("a:ext"))
    x = _emu_to_px(int(off.get("x", "0"))) if off is not None else 0.0
    y = _emu_to_px(int(off.get("y", "0"))) if off is not None else 0.0
    w = _emu_to_px(int(ext.get("cx", "0"))) if ext is not None else 0.0
    h = _emu_to_px(int(ext.get("cy", "0"))) if ext is not None else 0.0
    rot = int(xfrm.get("rot", "0")) / 60000.0
    return x, y, w, h, rot


def _find_layout_sp(layout_part, idx: int) -> Optional[etree._Element]:
    """Find the <p:sp> in a slideLayout that has <p:ph idx='idx'>."""
    sp_tree = layout_part._element.find(_x("p:cSld"))
    if sp_tree is None:
        return None
    sp_tree = sp_tree.find(_x("p:spTree"))
    if sp_tree is None:
        return None
    for sp in sp_tree.findall(_x("p:sp")):
        if _ph_idx(sp) == idx:
            return sp
    return None


def _find_master_sp(master_part, ph_type: Optional[str], idx: int) -> Optional[etree._Element]:
    """Find the <p:sp> in a slideMaster by ph type (preferred) or idx."""
    sp_tree = master_part._element.find(_x("p:cSld"))
    if sp_tree is None:
        return None
    sp_tree = sp_tree.find(_x("p:spTree"))
    if sp_tree is None:
        return None

    # First pass: match by type
    if ph_type is not None:
        for sp in sp_tree.findall(_x("p:sp")):
            if _ph_type(sp) == ph_type:
                return sp

    # Second pass: match by idx
    for sp in sp_tree.findall(_x("p:sp")):
        if _ph_idx(sp) == idx:
            return sp

    return None


def _get_master_tx_styles(master_part) -> Optional[etree._Element]:
    """Return the <p:txStyles> element from a slide master part."""
    return master_part._element.find(_x("p:txStyles"))


def _get_theme_el(master) -> Optional[etree._Element]:
    """Return the <a:theme> element from the theme part associated with a slide master."""
    try:
        part = master.part if hasattr(master, "part") else master
        theme_part = part.part_related_by(
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
        )
        if hasattr(theme_part, "_element"):
            return theme_part._element
        return etree.fromstring(theme_part.blob)
    except Exception:
        return None


# ---------------------------------------------------------------------------
# Picture helpers (P4)
# ---------------------------------------------------------------------------

# Microsoft Office's SVG-blip extension URI (drawingML 2016).
_SVG_BLIP_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"
_ASVG_NS = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"


def _resolve_blip_asset_ref(blip_fill_el: Optional[etree._Element], part) -> Optional[str]:
    """Resolve a ``<a:blipFill>``'s embedded image to a media basename.

    Reads ``<a:blip>`` from *blip_fill_el* and follows the relationship on
    *part* to find the target media partname. Returns the basename only
    (e.g. ``"image1.png"``) so it matches what
    :func:`pictures.extract.extract_pictures` writes to
    ``deck/public/assets/``.

    **SVG preference:** PowerPoint stores SVG icons as a PNG rasterization
    referenced by the standard ``r:embed`` attribute, *plus* an SVG original
    referenced from an ``<a:extLst><a:ext uri="{96DAC...}"><asvg:svgBlip
    r:embed="..."/>`` extension. When that extension is present we follow
    the SVG rId — vector graphics scale cleanly to slide canvas; the PNG
    fallback's resolution is fixed (often 96 DPI) so it looks pixelated on
    a 1280 / 1920-px slide.

    Returns ``None`` if no usable image relationship is found.
    """
    if blip_fill_el is None:
        return None
    blip = blip_fill_el.find(_x("a:blip"))
    if blip is None:
        return None

    # Prefer the SVG override (asvg:svgBlip) when present.
    ext_lst = blip.find(_x("a:extLst"))
    svg_rid: Optional[str] = None
    if ext_lst is not None:
        for ext in ext_lst.findall(_x("a:ext")):
            if ext.get("uri") == _SVG_BLIP_EXT_URI:
                svg_blip = ext.find(f"{{{_ASVG_NS}}}svgBlip")
                if svg_blip is not None:
                    svg_rid = svg_blip.get(_x("r:embed"))
                    break

    rid = svg_rid or blip.get(_x("r:embed"))
    if rid is None:
        # `r:link` (linked external image) is out of scope — we only handle embedded.
        return None
    try:
        rel = part.rels[rid]
    except KeyError:
        return None
    if "image" not in rel.reltype:
        return None
    try:
        partname = rel.target_part.partname
    except Exception:
        return None
    return Path(str(partname)).name


def _shape_clip_path(
    sp_pr: Optional[etree._Element],
    width_px: float,
    height_px: float,
) -> Optional[str]:
    """Build a CSS ``clip-path`` value for *sp_pr*'s geometry, or ``None``.

    Handles two PPT geometry forms on the shape's ``<p:spPr>``:

      - ``<a:prstGeom prst="..."/>``: maps to a preset clip-path / border-radius
        via :func:`pictures.geometry.preset_to_css`. ``rect`` (the default) and
        unknown presets return ``None`` (caller's wrapper renders as a plain
        rectangle).
      - ``<a:custGeom>``: freeform path; converted via
        :func:`pictures.geometry.cust_geom_to_clip_path`.

    Returns the ready-to-paste CSS value (e.g. ``'polygon(...)'`` or
    ``'path("M ... Z")'``) or ``None`` when no clipping should be applied.
    """
    if sp_pr is None:
        return None

    # custGeom first — if both are present (unusual but valid), the freeform
    # path is more specific.
    cust = sp_pr.find(_x("a:custGeom"))
    if cust is not None and width_px > 0 and height_px > 0:
        clip = cust_geom_to_clip_path(cust, width_px, height_px)
        if clip:
            return clip

    prst = sp_pr.find(_x("a:prstGeom"))
    if prst is not None:
        name = prst.get("prst")
        if name and name != "rect":
            av_lst = prst.find(_x("a:avLst"))
            avs: dict[str, int] = {}
            if av_lst is not None:
                for gd in av_lst.findall(_x("a:gd")):
                    gd_name = gd.get("name")
                    fmla = gd.get("fmla", "")
                    if gd_name and fmla.startswith("val "):
                        try:
                            avs[gd_name] = int(fmla[4:])
                        except ValueError:
                            continue
            geom = preset_to_css(
                name,
                int(width_px) if width_px else 0,
                int(height_px) if height_px else 0,
                avs or None,
            )
            if geom is not None:
                if geom.get("clip_path"):
                    return geom["clip_path"]
                # border-radius is handled by the wrapper style instead of
                # clip-path; returning None lets the emit layer skip clip-path
                # entirely and apply border-radius via its own pathway later.
                # For now placeholders don't propagate border-radius, so a
                # rounded-rect placeholder still renders as a plain rect.
                # Future work.
    return None


def _read_prst_geom(
    sp_pr: Optional[etree._Element],
) -> tuple[Optional[str], Optional[dict[str, int]]]:
    """Return (preset_name, adjust_values) for a shape's ``<a:prstGeom>``.

    Returns ``(None, None)`` when *sp_pr* is missing or has no prstGeom child,
    or when the preset is a freeform/custom geometry. Adjust values are read
    from ``<a:avLst><a:gd name="adj" fmla="val N"/></a:avLst>``; only the
    ``val N`` formula form is supported (the only one that maps directly to
    a CSS clip-path parameter — caller may pass the returned dict to
    :func:`pictures.geometry.preset_to_css`).
    """
    if sp_pr is None:
        return None, None
    prst = sp_pr.find(_x("a:prstGeom"))
    if prst is None:
        return None, None
    name = prst.get("prst")
    if not name:
        return None, None
    av_lst = prst.find(_x("a:avLst"))
    if av_lst is None:
        return name, None
    avs: dict[str, int] = {}
    for gd in av_lst.findall(_x("a:gd")):
        gd_name = gd.get("name")
        fmla = gd.get("fmla", "")
        if gd_name and fmla.startswith("val "):
            try:
                avs[gd_name] = int(fmla[4:])
            except ValueError:
                continue
    return name, (avs or None)


def _read_cnv_pr(parent_nv: Optional[etree._Element]) -> tuple[int, str]:
    """Return (shape_id, alt_text) from a ``<p:nvPicPr>`` or ``<p:nvSpPr>``.

    *parent_nv* should be the container that holds ``<p:cNvPr>``. Alt text is
    taken from ``@descr`` first, then ``@title``, then empty. Shape id defaults
    to 0 when not parsable.
    """
    if parent_nv is None:
        return 0, ""
    c_nv_pr = parent_nv.find(_x("p:cNvPr"))
    if c_nv_pr is None:
        return 0, ""
    try:
        shape_id = int(c_nv_pr.get("id", "0"))
    except ValueError:
        shape_id = 0
    alt_text = c_nv_pr.get("descr") or c_nv_pr.get("title") or ""
    return shape_id, alt_text


def _parse_pic(
    pic_el: etree._Element,
    slide_part,
    order_index: int,
) -> Optional[Picture]:
    """Parse a ``<p:pic>`` shape into a :class:`Picture`.

    A ``<p:pic>`` element is normally a free-floating picture, but OOXML
    also lets a ``<p:pic>`` carry a ``<p:ph>`` reference inside its
    ``<p:nvPicPr>/<p:nvPr>`` — that's how PowerPoint stores a picture
    placeholder whose default fill has been overridden with a specific
    image on the slide. Treating those as free pics caused a serious
    regression: the placeholder's idx never landed in ``slide_ph_idxes``,
    so the layout's matching idx got inherited and emitted as a SECOND
    picture next to the slide's version (the "icons shown twice" bug
    on the IU template's slides 19, 20, …).

    So: detect ``<p:ph>`` inside ``<p:nvPicPr>/<p:nvPr>``, and if present
    surface this Picture as ``is_placeholder=True`` with the placeholder
    idx. Layout inheritance then correctly skips the matching layout
    placeholder.

    Returns ``None`` only when the picture lacks both a resolvable image and
    geometry (degenerate case worth skipping silently).
    """
    nv_pic_pr = pic_el.find(_x("p:nvPicPr"))
    shape_id, alt_text = _read_cnv_pr(nv_pic_pr)

    # Placeholder detection — OOXML allows <p:pic> to be a picture
    # placeholder when it carries a <p:ph> in <p:nvPicPr>/<p:nvPr>.
    is_placeholder = False
    ph_idx: Optional[int] = None
    if nv_pic_pr is not None:
        nv_pr = nv_pic_pr.find(_x("p:nvPr"))
        if nv_pr is not None:
            ph_el = nv_pr.find(_x("p:ph"))
            if ph_el is not None:
                is_placeholder = True
                idx_attr = ph_el.get("idx")
                if idx_attr is not None:
                    try:
                        ph_idx = int(idx_attr)
                    except ValueError:
                        ph_idx = None

    blip_fill = pic_el.find(_x("p:blipFill"))
    asset_ref = _resolve_blip_asset_ref(blip_fill, slide_part)

    sp_pr = pic_el.find(_x("p:spPr"))
    x, y, w, h, rot_deg = _get_sp_position(pic_el)
    preset_name, preset_av = _read_prst_geom(sp_pr)
    effects = parse_effects(sp_pr, blip_fill)

    # Degenerate skip: no image, no size — nothing to render.
    if asset_ref is None and w == 0 and h == 0:
        return None

    return Picture(
        asset_ref=asset_ref,
        x_px=x,
        y_px=y,
        width_px=w,
        height_px=h,
        rotation_deg=rot_deg,
        preset_geom=preset_name,
        preset_geom_av=preset_av,
        effects=effects,
        alt_text=alt_text,
        shape_id=shape_id,
        is_placeholder=is_placeholder,
        ph_idx=ph_idx,
        order_index=order_index,
    )


def _parse_picture_placeholder(
    sp_el: etree._Element,
    slide_part,
    layout_part,
    layout_sp: Optional[etree._Element],
    idx: int,
    order_index: int,
) -> Picture:
    """Parse a picture-typed ``<p:sp>`` (``<p:ph type="pic">``) into a :class:`Picture`.

    blipFill cascade: slide-level ``<p:blipFill>`` wins; if missing or has no
    embed, fall back to the layout's blipFill on the matching layout shape
    (resolved against the layout part's rels).

    Geometry cascade: slide-level xfrm wins; if zero-sized, fall back to
    layout's xfrm.

    Returns a :class:`Picture` even when no image is bound on either level —
    in that case ``asset_ref`` is ``None`` and emit renders an empty
    positioned box (mirroring PPT's behaviour for un-bound picture
    placeholders).
    """
    nv_sp_pr = sp_el.find(_x("p:nvSpPr"))
    shape_id, alt_text = _read_cnv_pr(nv_sp_pr)

    sp_pr = sp_el.find(_x("p:spPr"))

    # blipFill cascade: slide → layout
    slide_blip_fill = sp_pr.find(_x("a:blipFill")) if sp_pr is not None else None
    asset_ref = _resolve_blip_asset_ref(slide_blip_fill, slide_part)
    effective_blip_fill = slide_blip_fill if asset_ref is not None else None
    if asset_ref is None and layout_sp is not None:
        layout_sp_pr = layout_sp.find(_x("p:spPr"))
        layout_blip_fill = (
            layout_sp_pr.find(_x("a:blipFill")) if layout_sp_pr is not None else None
        )
        layout_asset = _resolve_blip_asset_ref(layout_blip_fill, layout_part)
        if layout_asset is not None:
            asset_ref = layout_asset
            effective_blip_fill = layout_blip_fill

    # Geometry cascade — rotation cascades along with x/y/w/h.
    # When the slide-level shape omits its own xfrm (common when a slide
    # uses a placeholder verbatim from the layout), inherit the layout's
    # rotation too. Without this, a layout-defined tilt (e.g. the IU
    # template's slideLayout4 sets a -240000 EMU rot ≈ -4° on its image
    # placeholder idx=24) is silently dropped on the slide.
    x, y, w, h, rot_deg = _get_sp_position(sp_el)
    if w == 0 and h == 0 and layout_sp is not None:
        x, y, w, h, rot_deg = _get_sp_position(layout_sp)

    preset_name, preset_av = _read_prst_geom(sp_pr)
    effects = parse_effects(sp_pr, effective_blip_fill)

    return Picture(
        asset_ref=asset_ref,
        x_px=x,
        y_px=y,
        width_px=w,
        height_px=h,
        rotation_deg=rot_deg,
        preset_geom=preset_name,
        preset_geom_av=preset_av,
        effects=effects,
        alt_text=alt_text,
        shape_id=shape_id,
        is_placeholder=True,
        ph_idx=idx,
        order_index=order_index,
    )


def _parse_inherited_text_placeholder(
    src_sp: etree._Element,
    ph_type: Optional[str],
    idx: Optional[int],
    master,
    master_tx_styles,
    theme_el,
    clr_map,
    typefaces: set[str],
    order_index: int = 0,
) -> Optional[Placeholder]:
    """Parse a layout-only or master-only TEXT placeholder for a slide.

    Used when the slide itself doesn't declare a placeholder with this idx
    but the layout or master does (most commonly: footer / date /
    slidenum / body where the slide is "blank" and inherits from upstream).

    Geometry, fill, txBody all resolve against the source shape ONLY (no
    cross-cascade — there's no slide-level override to merge). Default-run
    / default-para still flow through the full cascade so font / color /
    bullet styling come from the master's <p:txStyles>.

    Returns ``None`` if the source shape's bbox is degenerate.
    """
    # Geometry
    x_px, y_px, w_px, h_px, rot_deg = _get_sp_position(src_sp)
    if w_px == 0 or h_px == 0:
        # If the layout shape has no own geometry, try master's matching shape.
        if master is not None and ph_type is not None:
            master_sp = _find_master_sp(master, ph_type, idx or -1)
            if master_sp is not None:
                x_px, y_px, w_px, h_px, rot_deg = _get_sp_position(master_sp)
    if w_px == 0 or h_px == 0:
        return None

    # Defaults via full cascade. slide_sp=src_sp because there's no real
    # slide-level shape; treat the inherited shape as the leaf of the chain.
    # include_slide_paragraph=False: src_sp is a layout/master shape, so its
    # first paragraph is prompt text, not authored content — reading it here
    # would reintroduce the slides-12/14 leak. Only lstStyle contributes.
    # ph_type defaults to "body" when the placeholder element omits @type
    # (mirrors the slide-level normalisation in parse()).
    effective_ph_type = ph_type if ph_type is not None else "body"
    default_run, default_para = resolve_placeholder(
        slide_sp=src_sp,
        layout_ph=None,
        master_ph=None,
        master_tx_styles=master_tx_styles,
        theme_el=theme_el,
        ph_type=effective_ph_type,
        level=0,
        clr_map=clr_map,
        include_slide_paragraph=False,
    )
    if default_run.font_family:
        typefaces.add(default_run.font_family)

    fill = _resolve_fill(src_sp, theme_el, clr_map, layout_sp=None)
    if fill is None:
        fill = NoFill()

    # Text frame from the inherited shape's txBody.
    tx_body = src_sp.find(_x("p:txBody"))
    text_frame = None
    if tx_body is not None and not _txbody_is_empty(tx_body):
        text_frame = _parse_text_frame(
            tx_body, default_run, default_para, theme_el, clr_map,
        )
        if text_frame is not None:
            for para in text_frame.paragraphs:
                for run in para.runs:
                    if run.font_family:
                        typefaces.add(run.font_family)

    shape_autofit, wrap_text = _resolve_body_pr_cascade(
        slide_sp=None, layout_sp=src_sp, master_sp=None,
    )
    return Placeholder(
        idx=idx if idx is not None else -1,
        type=ph_type,
        x_px=x_px,
        y_px=y_px,
        width_px=w_px,
        height_px=h_px,
        rotation_deg=rot_deg,
        fill=fill,
        opacity=1.0,
        text_frame=text_frame,
        default_run_props=default_run,
        default_para_props=default_para,
        is_prompt_fallback=False,
        shape_autofit=shape_autofit,
        wrap_text=wrap_text,
        order_index=order_index,
    )


def _parse_layout_only_picture_placeholder(
    layout_sp: etree._Element,
    layout_part,
    idx: int,
    order_index: int,
) -> Optional[Picture]:
    """Parse a layout-level ``<p:sp>`` picture placeholder NOT redeclared on the slide.

    OOXML semantics: a layout's picture-typed placeholder renders on every
    slide using that layout unless the slide overrides it by redeclaring the
    same idx. ``parse.py`` only walks the slide's own spTree, so without this
    helper layout-only picture placeholders never reach the resolved model
    (this was the cause of "16 assets extracted but only 5 rendered" against
    the IU template).

    Both blipFill and geometry resolve entirely against the layout part — no
    slide-level cascade since the slide has no matching shape.

    Returns ``None`` if the layout shape itself is degenerate (no image and
    zero geometry).
    """
    nv_sp_pr = layout_sp.find(_x("p:nvSpPr"))
    shape_id, alt_text = _read_cnv_pr(nv_sp_pr)

    sp_pr = layout_sp.find(_x("p:spPr"))
    blip_fill = sp_pr.find(_x("a:blipFill")) if sp_pr is not None else None
    asset_ref = _resolve_blip_asset_ref(blip_fill, layout_part)

    x, y, w, h, rot_deg = _get_sp_position(layout_sp)
    if asset_ref is None and w == 0 and h == 0:
        return None

    preset_name, preset_av = _read_prst_geom(sp_pr)
    effects = parse_effects(sp_pr, blip_fill if asset_ref is not None else None)

    return Picture(
        asset_ref=asset_ref,
        x_px=x,
        y_px=y,
        width_px=w,
        height_px=h,
        rotation_deg=rot_deg,
        preset_geom=preset_name,
        preset_geom_av=preset_av,
        effects=effects,
        alt_text=alt_text,
        shape_id=shape_id,
        is_placeholder=True,
        ph_idx=idx,
        order_index=order_index,
    )


# ---------------------------------------------------------------------------
# Main parse function
# ---------------------------------------------------------------------------

def parse(pptx_path: Path) -> Presentation:
    """Parse a PPTX file and return a fully-resolved Presentation model.

    Only text-bearing placeholders are extracted (v1 scope).
    """
    # Lazy import to break circular dependency: shapes.parse depends on
    # private helpers defined later in this module.
    from slidecraft.importer.shapes.parse import walk_text_shapes

    prs = pptx.Presentation(str(pptx_path))

    # Canvas dimensions
    sld_sz = prs._element.find(_x("p:sldSz"))
    if sld_sz is not None:
        canvas_w = int(int(sld_sz.get("cx", "9144000")) / _EMU_PER_PX)
        canvas_h = int(int(sld_sz.get("cy", "5143500")) / _EMU_PER_PX)
    else:
        canvas_w, canvas_h = 1920, 1080

    slides: list[Slide] = []
    typefaces: set[str] = set()

    for slide_idx, slide in enumerate(prs.slides, start=1):
        layout = slide.slide_layout
        master = layout.slide_master
        theme_el = _get_theme_el(master)
        master_tx_styles = _get_master_tx_styles(master)
        # Effective color map for THIS slide, honoring <p:clrMapOvr> on the slide
        # and its layout. A dark "title"/cover layout flips bg1<->tx1 via an
        # override; resolving only the master's <p:clrMap> would invert every
        # bg1/tx1 fill and font color on that cover (see get_effective_clr_map).
        clr_map = get_effective_clr_map(
            slide._element, layout._element, master._element
        )

        bg_fill = _resolve_background(slide, prs, theme_el, clr_map)

        placeholders: list[Placeholder] = []
        pictures: list[Picture] = []

        # Iterate shapes via lxml on the slide's XML
        cSld = slide._element.find(_x("p:cSld"))
        sp_tree = cSld.find(_x("p:spTree")) if cSld is not None else None
        if sp_tree is None:
            slides.append(Slide(index=slide_idx, placeholders=[], background_fill=bg_fill))
            continue

        # Record document order so picture/placeholder interleaving is reproducible
        # in emit. Each direct child of <p:spTree> gets a sequential index based
        # on its position in the underlying XML. We build the map by walking the
        # tree ONCE and using a dict keyed by the element wrapper itself; lxml
        # guarantees same-C-node ⇒ equal wrappers (==), so subsequent ``findall``
        # results can look up their position reliably via the helper below.
        sp_tree_children = list(sp_tree)
        def _doc_pos(el) -> int:
            # lxml's __eq__ compares the underlying C element identity, so
            # list.index() finds the same XML node even if findall() returns
            # a different Python proxy object. id()-keyed dicts are unreliable
            # because lxml's proxy lifetime is not tied to Python object lifetime.
            try:
                return sp_tree_children.index(el)
            except ValueError:
                return 0

        # ----- Picture placeholders + free <p:pic> -----
        # These are walked before the text-placeholder loop only because the
        # text-loop's `continue` path skips picture-typed phs; we still rely on
        # the same per-element traversal otherwise. Order is preserved via
        # the `order_index` field on Picture so emit can interleave correctly.
        for pic_el in sp_tree.findall(_x("p:pic")):
            picture = _parse_pic(
                pic_el,
                slide.part,
                _doc_pos(pic_el),
            )
            if picture is not None:
                pictures.append(picture)

        for sp_el in sp_tree.findall(_x("p:sp")):
            idx = _ph_idx(sp_el)
            if idx is None:
                continue  # Not a placeholder

            ph_type = _ph_type(sp_el)
            if ph_type == "pic":
                # Picture placeholder: cascades blipFill slide → layout
                layout_sp_for_pic = _find_layout_sp(layout, idx)
                picture = _parse_picture_placeholder(
                    sp_el,
                    slide.part,
                    layout.part,
                    layout_sp_for_pic,
                    idx,
                    _doc_pos(sp_el),
                )
                pictures.append(picture)
                continue
            if ph_type not in _TEXT_PH_TYPES:
                continue  # Non-text placeholder (chart, tbl, etc.) — still out of scope

            # Check if there's a txBody (for non-text types, skip gracefully)
            tx_body = sp_el.find(_x("p:txBody"))
            if tx_body is None:
                # Might still be valid as a text placeholder — only skip if type is definitely non-text
                if ph_type in ("pic", "chart", "tbl", "clipArt", "dgm", "media", "waveAudioFile"):
                    continue

            # Locate cascade elements
            layout_sp = _find_layout_sp(layout, idx)
            master_sp = _find_master_sp(master, ph_type, idx)

            # OOXML default: when <p:ph> omits the @type attribute (common for
            # "content" placeholders), the placeholder defaults to body
            # semantics. Without this normalisation, _txstyles_defaults falls
            # through to <p:otherStyle> and the master's <p:bodyStyle> bullet
            # character + indent never apply. Slide 10's content placeholder
            # (<p:ph idx="1"/>) hit this bug — body text rendered without the
            # "-" bullet marker the master defines.
            # Non-placeholder text shapes (in shapes/parse.py) still pass
            # ph_type=None explicitly to mean "use otherStyle" — that path is
            # unaffected because this normalisation only runs in the
            # placeholder loop.
            effective_ph_type = ph_type if ph_type is not None else "body"

            # Resolve defaults (cascade level 1–5)
            default_run, default_para = resolve_placeholder(
                slide_sp=sp_el,
                layout_ph=layout_sp,
                master_ph=master_sp,
                master_tx_styles=master_tx_styles,
                theme_el=theme_el,
                ph_type=effective_ph_type,
                level=0,
                clr_map=clr_map,
            )

            # Collect typefaces from defaults
            if default_run.font_family:
                typefaces.add(default_run.font_family)

            # Geometry — prefer slide-level xfrm, fall back to layout, then master
            x_px, y_px, w_px, h_px, rot_deg = _get_sp_position(sp_el)
            if w_px == 0 and h_px == 0:
                if layout_sp is not None:
                    x_px, y_px, w_px, h_px, rot_deg = _get_sp_position(layout_sp)
            if w_px == 0 and h_px == 0:
                if master_sp is not None:
                    x_px, y_px, w_px, h_px, rot_deg = _get_sp_position(master_sp)

            # Fill — cascades slide → layout
            fill = _resolve_fill(sp_el, theme_el, clr_map, layout_sp=layout_sp)

            # Opacity (not commonly set on placeholders, default 1.0)
            opacity = 1.0

            # Per-level cascade resolution (RC2). Resolve defaults for
            # every list-indent level 0..8 so that paragraphs at lvl=N can
            # inherit master's lvl(N+1)pPr defRPr (font_size, font_family,
            # bold, etc.) when their slide-level <a:rPr> doesn't override.
            # Without this, tmp2 slide 10's lvl=3 sub-bullets rendered at
            # the placeholder's level-0 font-size (32pt) instead of
            # master's lvl4pPr 20pt — visibly LARGER than the level-0
            # bullets they were nested under.
            per_level_defaults: dict[int, tuple[Run, Paragraph]] = {}
            for lvl_n in range(9):
                if lvl_n == 0:
                    per_level_defaults[0] = (default_run, default_para)
                else:
                    # include_slide_paragraph=False: the slide's first
                    # paragraph is level-0 content — its explicit pPr/rPr
                    # must not bleed into level-N defaults, which come from
                    # the per-level lstStyle cascade alone.
                    lvl_run, lvl_para = resolve_placeholder(
                        slide_sp=sp_el,
                        layout_ph=layout_sp,
                        master_ph=master_sp,
                        master_tx_styles=master_tx_styles,
                        theme_el=theme_el,
                        ph_type=effective_ph_type,
                        level=lvl_n,
                        clr_map=clr_map,
                        include_slide_paragraph=False,
                    )
                    per_level_defaults[lvl_n] = (lvl_run, lvl_para)

            # TextFrame — with prompt-fallback for empty slide-level content
            text_frame = None
            is_prompt_fallback = False
            if tx_body is not None:
                if (
                    _txbody_is_empty(tx_body)
                    and layout_sp is not None
                    and _layout_ph_has_custom_prompt(layout_sp)
                ):
                    # Slide placeholder is empty; use layout's prompt text as the content.
                    # Only the text RUNS come from the layout; styling cascade is already
                    # resolved from slide → layout → master defaults (default_run/default_para).
                    layout_tx_body = layout_sp.find(_x("p:txBody"))
                    if layout_tx_body is not None and not _txbody_is_empty(layout_tx_body):
                        text_frame = _parse_text_frame(
                            layout_tx_body, default_run, default_para, theme_el, clr_map,
                            per_level_defaults=per_level_defaults,
                        )
                        is_prompt_fallback = True
                if text_frame is None:
                    layout_tx_body = (
                        layout_sp.find(_x("p:txBody")) if layout_sp is not None else None
                    )
                    master_tx_body = (
                        master_sp.find(_x("p:txBody")) if master_sp is not None else None
                    )
                    text_frame = _parse_text_frame(
                        tx_body, default_run, default_para, theme_el, clr_map,
                        layout_tx_body=layout_tx_body,
                        master_tx_body=master_tx_body,
                        per_level_defaults=per_level_defaults,
                    )
                # Collect typefaces from actual runs
                if text_frame is not None:
                    for para in text_frame.paragraphs:
                        for run in para.runs:
                            if run.font_family:
                                typefaces.add(run.font_family)

            # Geometry shape (prstGeom / custGeom) cascade: slide → layout
            # → master. The drawer-shape chip on the IU title page is a
            # <a:custGeom> on the LAYOUT placeholder; without this cascade
            # the slide-level wrapper falls back to a plain rect.
            clip_path: Optional[str] = None
            for src_sp in (sp_el, layout_sp, master_sp):
                if src_sp is None:
                    continue
                src_sp_pr = src_sp.find(_x("p:spPr"))
                clip_path = _shape_clip_path(src_sp_pr, w_px, h_px)
                if clip_path is not None:
                    break

            shape_autofit, wrap_text = _resolve_body_pr_cascade(
                slide_sp=sp_el, layout_sp=layout_sp, master_sp=master_sp,
            )
            placeholders.append(Placeholder(
                idx=idx,
                type=ph_type,
                x_px=x_px,
                y_px=y_px,
                width_px=w_px,
                height_px=h_px,
                rotation_deg=rot_deg,
                fill=fill,
                opacity=opacity,
                text_frame=text_frame,
                default_run_props=default_run,
                default_para_props=default_para,
                is_prompt_fallback=is_prompt_fallback,
                clip_path=clip_path,
                shape_autofit=shape_autofit,
                wrap_text=wrap_text,
                order_index=_doc_pos(sp_el),
            ))

        # Surface layout-only picture placeholders. OOXML: a layout-level
        # <p:sp type="pic"> renders on every slide using that layout unless
        # the slide redeclares the same idx. Without this pass, the IU
        # template's layout-bound decorative images (logos, sidebars, sample
        # photos baked into the layout's blipFill) never reach the model.
        slide_ph_idxes = {p.idx for p in placeholders} | {
            p.ph_idx for p in pictures if p.is_placeholder and p.ph_idx is not None
        }
        layout_sp_tree = (
            layout._element.find(_x("p:cSld") + "/" + _x("p:spTree"))
            if layout._element is not None
            else None
        )
        if layout_sp_tree is None:
            # Fallback: find with explicit walk
            l_csld = layout._element.find(_x("p:cSld"))
            layout_sp_tree = l_csld.find(_x("p:spTree")) if l_csld is not None else None
        if layout_sp_tree is not None:
            # Layout-inherited shapes get NEGATIVE order_index so they always
            # render BEHIND anything the slide itself authored. Without this,
            # an inherited background placeholder (e.g. picture-22 covering
            # the full slide) would emit AFTER slide-level title placeholders
            # in DOM order — putting it on top and visually hiding the title.
            # OOXML semantics treat the layout as background; we mirror that
            # with a negative base so sorting by order_index ASCENDING puts
            # all layout content first (= earlier in DOM = lower z-index).
            base = -10_000
            for i, lsp in enumerate(layout_sp_tree.findall(_x("p:sp"))):
                lidx = _ph_idx(lsp)
                if lidx is None:
                    continue
                if lidx in slide_ph_idxes:
                    continue  # Slide overrides this layout placeholder.
                if _ph_type(lsp) != "pic":
                    continue
                inherited = _parse_layout_only_picture_placeholder(
                    lsp, layout.part, lidx, base + i,
                )
                if inherited is not None:
                    pictures.append(inherited)

            # Also walk layout-level free <p:pic> shapes (e.g. the IU template's
            # decorative logo lives as a <p:pic> directly on the layout, not
            # inside a placeholder). These render under every slide that uses
            # the layout, so we surface them per-slide.
            for j, lpic in enumerate(layout_sp_tree.findall(_x("p:pic"))):
                inherited_free = _parse_pic(
                    lpic, layout.part, base + 5_000 + j,
                )
                if inherited_free is not None:
                    pictures.append(inherited_free)

            # Surface layout-only TEXT placeholders. OOXML: a layout-level
            # <p:sp> with <p:ph> renders on every slide using that layout
            # unless the slide redeclares the same idx. Vorlage_Foliensatz
            # slide 5 (and many others) doesn't redeclare body / title /
            # footer placeholders that the layout carries — those simply
            # inherit. Without this pass, those slides come out near-empty.
            #
            # Dedup by (type, idx). The slide can use a different idx for
            # the same conceptual placeholder (IU template's footer is
            # idx=11 on the slide but idx=3 on the master); a per-type
            # set catches those so we don't double-add. Singleton types
            # (title / ctrTitle / dt / ftr / sldNum) are unique per slide
            # — any presence on slide counts as "already declared".
            _SINGLETON_PH_TYPES = frozenset({
                "title", "ctrTitle", "dt", "ftr", "sldNum",
            })
            slide_ph_types = {p.type for p in placeholders if p.type is not None}

            for lsp_i, lsp in enumerate(layout_sp_tree.findall(_x("p:sp"))):
                lidx = _ph_idx(lsp)
                lph_type = _ph_type(lsp)
                if lph_type == "pic":
                    continue
                if lidx in slide_ph_idxes:
                    continue
                # Singleton types: skip if slide has any placeholder of this type.
                if lph_type in _SINGLETON_PH_TYPES and lph_type in slide_ph_types:
                    continue
                if lidx is None and lph_type is None:
                    continue
                if lph_type not in _TEXT_PH_TYPES and lph_type is not None:
                    continue
                # Layout-inherited text placeholders go behind slide content
                # (same negative-base convention as inherited pictures, with
                # a small offset so they don't collide with the picture range).
                inherited_ph = _parse_inherited_text_placeholder(
                    lsp, lph_type, lidx, master,
                    master_tx_styles, theme_el, clr_map, typefaces,
                    order_index=base + 8_000 + lsp_i,
                )
                if inherited_ph is not None:
                    placeholders.append(inherited_ph)
                    if lidx is not None:
                        slide_ph_idxes.add(lidx)
                    if lph_type is not None:
                        slide_ph_types.add(lph_type)

        # Master-only TEXT placeholder inheritance INTENTIONALLY SKIPPED.
        #
        # PPT's master placeholders for footer / date / slide-number are
        # rendered on a slide only when `<p:hf ftr="1"/>` (etc.) opts in,
        # OR when the layout / slide itself redeclares the placeholder.
        # Without that opt-in, PPT does NOT show the master's content —
        # it stays purely as a styling template / positioning hint.
        #
        # An earlier version of this code unconditionally inherited
        # master ftr/dt/sldNum, which produced Quelle Text on every
        # tmp2 slide (Vorlage master's ftr placeholder holds the literal
        # text "Quelle Text: Autorennachname..."). That over-emitted on
        # divider slides and also double-emitted on slides 34-58/82 which
        # have their own slide-level Quelle.
        #
        # The layout walk above already inherits any layout-level
        # placeholders the user wants visible (e.g. body, title, dt with
        # actual content). Master content stays in its proper role as a
        # cascade source, not a rendered shape.
        pass

        # Master-level inheritance. OOXML: master shapes render under every
        # slide unless the slide or its layout sets showMasterSp="0". The IU
        # template's "black on light background" logo lives on the master and
        # only appears on slides whose layout doesn't override it with the
        # white variant. Dedup by position so a layout's own logo (e.g. the
        # white variant at the same top-right position) suppresses the
        # master's logo for that slide.
        slide_show_master = slide._element.get("showMasterSp")
        layout_show_master = (
            layout._element.get("showMasterSp")
            if layout._element is not None
            else None
        )
        effective_show_master = slide_show_master or layout_show_master or "1"
        if effective_show_master != "0":
            master_sp_tree = None
            m_csld = master._element.find(_x("p:cSld"))
            if m_csld is not None:
                master_sp_tree = m_csld.find(_x("p:spTree"))
            if master_sp_tree is not None:
                # Build a position set from already-collected pictures so
                # master pics overlapping a layout pic at the same xy can be
                # skipped (layout overrides). Rounded-int positions tolerate
                # subpixel EMU conversion noise.
                covered_positions = {
                    (round(p.x_px), round(p.y_px)) for p in pictures
                }
                # Master shapes go BEHIND layout shapes (which are themselves
                # behind slide content). More negative = further back = emitted
                # earlier in DOM = lower z-index.
                master_base = -20_000
                for k, mpic in enumerate(master_sp_tree.findall(_x("p:pic"))):
                    inherited_master = _parse_pic(
                        mpic, master.part, master_base + k,
                    )
                    if inherited_master is None:
                        continue
                    pos_key = (
                        round(inherited_master.x_px),
                        round(inherited_master.y_px),
                    )
                    if pos_key in covered_positions:
                        continue  # Layout's own pic at this xy overrides.
                    pictures.append(inherited_master)

        # Stable document order across both <p:pic> shapes and <p:ph type="pic">
        # placeholders (which arrive from four sources: slide-level <p:pic>,
        # slide-level <p:sp type="pic">, slide layout-inherited shapes, and
        # master-inherited <p:pic> shapes).
        pictures.sort(key=lambda p: p.order_index)

        # Layer 3 — non-placeholder text-bearing shapes from slide, layout,
        # and master spTrees. walk_text_shapes() honours showMasterSp on both
        # the slide and the layout, filters empty / zero-size shapes, and
        # resolves text defaults via inheritance.resolve_placeholder() with
        # ph_type=None (→ master <a:otherStyle>).
        text_shapes = walk_text_shapes(
            slide.part,
            master_tx_styles,
            theme_el,
            clr_map,
        )

        # Layer 4 — tables (<p:graphicFrame> / <a:tbl>) on the slide. Lazy
        # import to break the parse-shapes import cycle (tables.parse pulls
        # the same private helpers from this module).
        from slidecraft.importer.tables.parse import walk_tables
        tables = walk_tables(
            slide.part,
            master_tx_styles,
            theme_el,
            clr_map,
        )
        # Surface typefaces from cell content so fonts.css carries them.
        for tbl in tables:
            for row in tbl.cells:
                for cell in row:
                    if cell.default_run.font_family:
                        typefaces.add(cell.default_run.font_family)
                    if cell.text_frame is not None:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font_family:
                                    typefaces.add(run.font_family)
        # Surface typefaces used by these shapes so fonts.css carries them.
        for ts in text_shapes:
            if ts.default_run.font_family:
                typefaces.add(ts.default_run.font_family)
            if ts.text_frame is not None:
                for para in ts.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font_family:
                            typefaces.add(run.font_family)

        slides.append(Slide(
            index=slide_idx,
            placeholders=placeholders,
            background_fill=bg_fill,
            pictures=pictures,
            text_shapes=text_shapes,
            tables=tables,
        ))

    return Presentation(
        slides=slides,
        canvas_width_px=canvas_w,
        canvas_height_px=canvas_h,
        typefaces_referenced=typefaces,
    )


